from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

def add_title_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(45, 62, 80)  # Dark blue background
    
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text
    title_shape.text_frame.paragraphs[0].font.size = Pt(44)
    title_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        subtitle_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(236, 240, 241)  # Light gray text
        subtitle_shape.text_frame.paragraphs[0].font.size = Pt(28)
        subtitle_shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def add_section_slide(prs, title):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Using the blank layout
    # Set background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(52, 152, 219)  # Blue background
    
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    
    shape = slide.shapes.add_textbox(left, top, width, height)
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, content, include_icon=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    # Set light background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(236, 240, 241)  # Light gray background
    
    # Add title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(36)
    title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(44, 62, 80)  # Dark blue text
    
    # Add content
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for item in content:
        p = tf.add_paragraph()
        if isinstance(item, tuple):
            p.text = "• " + item[0] if item[1] > 0 else item[0]
            p.level = item[1]
        else:
            p.text = item
            p.level = 0
        
        p.font.size = Pt(24 if p.level == 0 else 20)
        p.font.color.rgb = RGBColor(44, 62, 80)  # Dark blue text
        p.space_after = Pt(12)  # Add space between items

def create_presentation():
    prs = Presentation()
    
    # Title slide
    add_title_slide(prs, "ReadRadar", "Your Personal Book Discovery Platform")
    
    # What is ReadRadar?
    add_section_slide(prs, "What is ReadRadar?")
    
    add_content_slide(prs, "Your Digital Reading Companion", [
        "📚 Discover and Track Your Reading Journey",
        ("Find your next favorite book", 1),
        ("Keep track of what you read", 1),
        ("Share your thoughts with others", 1),
        "🌟 Connect with Book Lovers",
        ("Join a community of readers", 1),
        ("Share reading lists", 1),
        ("Discover recommendations", 1)
    ])
    
    # Current Features
    add_section_slide(prs, "What Can You Do Now?")
    
    add_content_slide(prs, "Current Features", [
        "👤 Personal Account",
        ("Create your reading profile", 1),
        ("Secure login & data protection", 1),
        "📖 Book Management",
        ("Add books to your collection", 1),
        ("Write and read reviews", 1),
        ("Rate books you've read", 1),
        "🔍 Easy Navigation",
        ("Search for books", 1),
        ("Browse by genre", 1),
        ("View detailed book info", 1)
    ])
    
    # Coming Soon
    add_section_slide(prs, "What's Coming Next?")
    
    add_content_slide(prs, "Exciting New Features", [
        "🎯 Smart Search",
        ("Find books instantly", 1),
        ("Filter by multiple criteria", 1),
        ("Get personalized suggestions", 1),
        "🤝 Social Features",
        ("Follow other readers", 1),
        ("Share your reading lists", 1),
        ("Discuss books with friends", 1),
        "🎁 More Surprises",
        ("Dark mode for night reading", 1),
        ("Reading challenges & rewards", 1)
    ])
    
    # Timeline
    add_section_slide(prs, "When to Expect More?")
    
    add_content_slide(prs, "Development Timeline", [
        "📅 Weeks 1-2",
        ("Advanced search features", 1),
        ("Improved book discovery", 1),
        "📅 Weeks 3-4",
        ("Social features rollout", 1),
        ("Reading list enhancements", 1),
        "📅 Weeks 5-6",
        ("Performance optimization", 1),
        ("Final testing & polish", 1)
    ])
    
    # Thank You slide
    add_title_slide(prs, "Join Us on the Journey!", "ReadRadar - Where Every Book Finds Its Reader")
    
    # Save the presentation
    prs.save('docs/presentation/ReadRadar_Midterm.pptx')

if __name__ == '__main__':
    create_presentation()
